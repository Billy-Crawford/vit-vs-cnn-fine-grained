import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# PRÉSENTATION OFFICIELLE (22 SLIDES) — MASTER 1 IA — ViT vs CNN (CUB-200-2011)
# Palette stricte 3 couleurs : Navy sombre (#0F172A), Bleu accent (#2563EB), Fond doux (#F8FAFC)
# ==============================================================================

C_NAVY_DARK    = RGBColor(15, 23, 42)     # 1. Sombre principal (Navy / Slate-900)
C_BLUE_ACCENT  = RGBColor(37, 99, 235)    # 2. Accent principal (Bleu royal / Blue-600)
C_BG_PAGE      = RGBColor(248, 250, 252)  # 3. Fond neutre très doux (Slate-50)

# Couleurs neutres de structure
C_WHITE        = RGBColor(255, 255, 255)  # Blanc pur (cartes & contrastes)
C_BORDER_LIGHT = RGBColor(226, 232, 240)  # Bordures discrètes (Slate-200)
C_TEXT_MUTED   = RGBColor(100, 116, 139)  # Texte secondaire (Slate-500)

FONT_MAIN = "Arial"
TOTAL_SLIDES = 22

def create_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    REPO = os.path.dirname(os.path.abspath(__file__))

    # ── Helpers de dessin épurés ─────────────────────────────────
    def set_solid(shape, color=C_BLUE_ACCENT):
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def draw_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        set_solid(bg, C_BG_PAGE)
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
        set_solid(top_bar, C_BLUE_ACCENT)

    def draw_card(slide, l, t, w, h, bg_color=C_WHITE, border_color=C_BORDER_LIGHT, border_pt=1.0):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(border_pt)
        return card

    def add_header(slide, title):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.52), Inches(11.7), Inches(0.65))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(23)
        p_t.font.bold = True
        p_t.font.color.rgb = C_NAVY_DARK
        p_t.font.name = FONT_MAIN
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.22), Inches(2.6), Inches(0.03))
        set_solid(line, C_BLUE_ACCENT)

    def add_footer(slide, current_slide, total_slides=TOTAL_SLIDES):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.733), Inches(0.3))
        tf = footer_box.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"Master 1 IA - ViT vs CNN (CUB-200-2011)                                                                                            {current_slide}/{total_slides}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = C_TEXT_MUTED
        p.font.name = FONT_MAIN

    def add_pill_badge(slide, l, t, w, h, text, bg_color=C_BLUE_ACCENT, text_color=C_WHITE, font_size=10.0):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        set_solid(pill, bg_color)
        tf = pill.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_MAIN
        return pill

    def add_diamond_badge(slide, l, t, size, text, bg_color=C_BLUE_ACCENT, text_color=C_WHITE):
        d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, size, size)
        set_solid(d, bg_color)
        tf = d.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_MAIN
        return d

    # =========================================================================
    # SLIDE 1 : TITRE & ÉQUIPE
    # =========================================================================
    s1 = prs.slides.add_slide(blank)
    draw_bg(s1)

    r_panel = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), 0, Inches(5.033), prs.slide_height)
    set_solid(r_panel, C_WHITE)
    r_border = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), 0, Inches(0.02), prs.slide_height)
    set_solid(r_border, C_BORDER_LIGHT)

    tb_title = s1.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(7.0), Inches(2.8))
    tf1 = tb_title.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Vision Transformer"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_DARK
    
    p2 = tf1.add_paragraph()
    p2.text = "vs  CNN (ResNet-50)"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = C_BLUE_ACCENT
    
    p3 = tf1.add_paragraph()
    p3.text = "Étude comparative sur la classification fine-grained d'oiseaux"
    p3.font.size = Pt(16)
    p3.font.color.rgb = C_NAVY_DARK
    p3.space_before = Pt(8)
    p3.line_spacing = 1.4

    p4 = tf1.add_paragraph()
    p4.text = "Dataset CUB-200-2011  ·  200 espèces  ·  11 788 images"
    p4.font.size = Pt(12)
    p4.font.color.rgb = C_TEXT_MUTED
    p4.space_before = Pt(4)

    draw_card(s1, Inches(0.9), Inches(4.5), Inches(7.0), Inches(2.2), C_WHITE, C_BORDER_LIGHT)
    tb_stat = s1.shapes.add_textbox(Inches(1.1), Inches(4.65), Inches(6.6), Inches(1.9))
    tfs = tb_stat.text_frame
    tfs.word_wrap = True
    p = tfs.paragraphs[0]
    p.text = "POINTS FORTS DU PROJET"
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = C_BLUE_ACCENT
    p.space_after = Pt(6)

    p = tfs.add_paragraph()
    p.text = "• 12 configurations expérimentales complètes entraînées et évaluées.\n• Étude d'ablation rigoureuse (taille de patch, pré-entraînement, volume de données).\n• Application web interactive avec extraction réelle de carte d'attention ViT.\n• Déploiement multi-plateforme 100% dockerisé (FastAPI + Next.js)."
    p.font.size = Pt(11.5)
    p.font.color.rgb = C_NAVY_DARK
    p.line_spacing = 1.4

    team_title = s1.shapes.add_textbox(Inches(8.6), Inches(0.8), Inches(4.4), Inches(0.4))
    team_title.text_frame.paragraphs[0].text = "MEMBRES DE L'ÉQUIPE"
    team_title.text_frame.paragraphs[0].font.bold = True
    team_title.text_frame.paragraphs[0].font.size = Pt(13)
    team_title.text_frame.paragraphs[0].font.color.rgb = C_NAVY_DARK

    team_members = [
        ("RÔLE A — DATA / EXPERIMENT", "SONHOUIN Abdoul-raouf", "Pipeline data, EDA, splits stratifiés, augmentations PyTorch, préparation CUB."),
        ("RÔLE B — MODEL / RESEARCH", "NGARTOBAYE OUMAROU BILLY", "Implémentation ViT scratch, intégration timm/ResNet-50, 12 runs & MLflow."),
        ("RÔLE C — REPORTING / BACKEND", "YEYE Koffi Gagnon", "Backend FastAPI, Frontend Next.js, extraction d'attention ViT, Docker & rapport."),
    ]

    for i, (role, name, desc) in enumerate(team_members):
        top = Inches(1.4 + i * 1.8)
        draw_card(s1, Inches(8.6), top, Inches(4.4), Inches(1.6), C_WHITE, C_BORDER_LIGHT)
        rect_tag = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.6), top, Inches(0.08), Inches(1.6))
        set_solid(rect_tag, C_BLUE_ACCENT)
        
        tb = s1.shapes.add_textbox(Inches(8.85), top + Inches(0.12), Inches(4.0), Inches(1.35))
        tf = tb.text_frame; tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = role; p.font.bold = True; p.font.size = Pt(9.5); p.font.color.rgb = C_BLUE_ACCENT
        
        p2 = tf.add_paragraph()
        p2.text = name; p2.font.bold = True; p2.font.size = Pt(13); p2.font.color.rgb = C_NAVY_DARK; p2.space_after = Pt(2)
        
        p3 = tf.add_paragraph()
        p3.text = desc; p3.font.size = Pt(10); p3.font.color.rgb = C_TEXT_MUTED; p3.line_spacing = 1.35

    # =========================================================================
    # SLIDE 2 : PLAN DE LA SOUTENANCE (EXÉCUTIF & AÉRÉ)
    # =========================================================================
    s2 = prs.slides.add_slide(blank)
    draw_bg(s2)
    add_header(s2, "Plan de la Présentation")
    add_footer(s2, 2)

    sections_plan = [
        ("01", "Contexte, Problématique & Fondements Théoriques", "Défi fine-grained, dataset CUB-200 et comparaison CNN vs Vision Transformer"),
        ("02", "Architecture Globale & Organisation du Code", "Vue d'ensemble end-to-end et structure modulaire des répertoires"),
        ("03", "Pipeline Données & Analyse Exploratoire (Rôle A)", "Splits stratifiés (10% à 100%), résolutions hétérogènes & augmentations"),
        ("04", "Modélisation, Protocole & Suivi MLflow (Rôle B)", "ResNet-50, ViT timm, ViT custom scratch et traçabilité des 12 runs"),
        ("05", "Résultats Expérimentaux & Études d'Ablation", "Matrice comparative, sensibilité Data-Hungry et impact du pré-entraînement"),
        ("06", "Déploiement Fullstack, Attention ViT & Bilan (Rôle C)", "FastAPI, Next.js, hook d'attention, Docker Compose, limites & perspectives"),
    ]

    for idx, (num, titre, sous_titre) in enumerate(sections_plan):
        row = idx // 2
        col = idx % 2
        l = Inches(0.9 + col * 5.9)
        t = Inches(1.6 + row * 1.7)
        
        draw_card(s2, l, t, Inches(5.6), Inches(1.45), C_WHITE, C_BORDER_LIGHT)
        add_diamond_badge(s2, l + Inches(0.2), t + Inches(0.2), Inches(0.7), num, C_BLUE_ACCENT, C_WHITE)
        
        tb = s2.shapes.add_textbox(l + Inches(1.1), t + Inches(0.15), Inches(4.3), Inches(1.15))
        tf = tb.text_frame; tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = titre; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = sous_titre; p2.font.size = Pt(10); p2.font.color.rgb = C_TEXT_MUTED; p2.line_spacing = 1.35

    # =========================================================================
    # SLIDE 3 : CONTEXTE & DÉFI DU FINE-GRAINED
    # =========================================================================
    s3 = prs.slides.add_slide(blank)
    draw_bg(s3)
    add_header(s3, "1. Contexte & Défi de la Classification Fine-Grained")
    add_footer(s3, 3)

    cards_s3 = [
        ("Défi de la Classification Fine-Grained", [
            ("Variabilité inter-classe minime :", "\n  Les différences entre espèces sont extrêmement subtiles (nuances de plumage, courbure du bec, motifs des rémiges)."),
            ("Variabilité intra-classe forte :", "\n  Une même espèce présente des apparences très divergentes selon l'âge, le sexe, la saison, la pose et l'éclairage."),
            ("Complexité du contexte visuel :", "\n  Arrière-plans naturels denses (branches, feuillages) créant du bruit visuel perturbateur.")
        ]),
        ("Enjeu Scientifique & Algorithmique", [
            ("Localisation de détails discriminants :", "\n  Le modèle doit focaliser son apprentissage sur des micro-régions anatomiques clés tout en ignorant le décor."),
            ("Volume de données modeste :", "\n  Environ 30 images d'entraînement par espèce (5 094 images train au total), ce qui constitue un défi majeur pour les architectures sans a priori spatial.")
        ])
    ]

    for i, (titre, points) in enumerate(cards_s3):
        l = Inches(0.9 + i * 5.9)
        draw_card(s3, l, Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
        h_s3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.55), Inches(5.6), Inches(0.55))
        set_solid(h_s3, C_NAVY_DARK if i == 0 else C_BLUE_ACCENT)
        
        tb_h = s3.shapes.add_textbox(l + Inches(0.15), Inches(1.62), Inches(5.3), Inches(0.4))
        tb_h.text_frame.paragraphs[0].text = titre
        tb_h.text_frame.paragraphs[0].font.bold = True
        tb_h.text_frame.paragraphs[0].font.size = Pt(12)
        tb_h.text_frame.paragraphs[0].font.color.rgb = C_WHITE

        tb_b = s3.shapes.add_textbox(l + Inches(0.2), Inches(2.25), Inches(5.2), Inches(4.3))
        tf_b = tb_b.text_frame; tf_b.word_wrap = True
        for j, (b_t, r_t) in enumerate(points):
            p = tf_b.paragraphs[0] if j == 0 else tf_b.add_paragraph()
            run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
            run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(11); run2.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(8)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 4 : DATASET CUB-200-2011 & OBJECTIFS DE L'ÉTUDE
    # =========================================================================
    s4 = prs.slides.add_slide(blank)
    draw_bg(s4)
    add_header(s4, "2. Le Benchmark CUB-200-2011 & Objectifs de l'Étude")
    add_footer(s4, 4)

    draw_card(s4, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s4, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "BENCHMARK DE RÉFÉRENCE", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s4_l = s4.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s4_l = tb_s4_l.text_frame; tf_s4_l.word_wrap = True

    pts_cub = [
        ("• Origine & Structure :", "\n  Caltech-UCSD Birds-200-2011, benchmark de référence internationale pour la vision par ordinateur fine-grained."),
        ("• Volumétrie Totale :", "\n  11 788 images réparties sur 200 espèces d'oiseaux nord-américains (41 à 60 images par classe)."),
        ("• Split Officiel :", "\n  5 994 images d'entraînement / 5 794 images de test (évaluation finale standardisée).")
    ]
    for b_t, r_t in pts_cub:
        p = tf_s4_l.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(11); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.35

    draw_card(s4, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s4, Inches(7.05), Inches(1.8), Inches(3.2), Inches(0.32), "OBJECTIFS SCIENTIFIQUES", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s4_r = s4.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s4_r = tb_s4_r.text_frame; tf_s4_r.word_wrap = True

    pts_obj = [
        ("1. Confrontation Objective CNN vs ViT :", "\n  Mesurer rigoureusement les performances relatives de ResNet-50 et du Vision Transformer sur ce benchmark exigeant."),
        ("2. Validation Empirique de l'Hypothèse Data-Hungry :", "\n  Quantifier comment les deux paradigmes réagissent lorsque l'on réduit drastiquement les données (10%, 25%, 50%, 100%)."),
        ("3. Restitution Applicative & Démonstrateur :", "\n  Fournir une plateforme web complète avec explication visuelle temps réel par carte d'attention.")
    ]
    for b_t, r_t in pts_obj:
        p = tf_s4_r.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(11); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 5 : FONDEMENTS THÉORIQUES — CNN vs VISION TRANSFORMER
    # =========================================================================
    s5 = prs.slides.add_slide(blank)
    draw_bg(s5)
    add_header(s5, "3. Fondements Théoriques : CNN vs Vision Transformer")
    add_footer(s5, 5)

    draw_card(s5, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    h_cnn = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.55), Inches(5.6), Inches(0.6))
    set_solid(h_cnn, C_NAVY_DARK)
    tb_cnn_h = s5.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.2), Inches(0.4))
    tb_cnn_h.text_frame.paragraphs[0].text = "Réseaux Convolutifs — ResNet-50 (CNN)"
    tb_cnn_h.text_frame.paragraphs[0].font.bold = True; tb_cnn_h.text_frame.paragraphs[0].font.size = Pt(13); tb_cnn_h.text_frame.paragraphs[0].font.color.rgb = C_WHITE

    tb_cnn_b = s5.shapes.add_textbox(Inches(1.15), Inches(2.35), Inches(5.1), Inches(4.2))
    tf_cb = tb_cnn_b.text_frame; tf_cb.word_wrap = True
    cnn_pts = [
        ("• Biais Inductif de Localité :", "\n  Les filtres de convolution traitent des voisinages de pixels réduits (ex. 3×3, 7×7)."),
        ("• Invariance par Translation :", "\n  Reconnaissance d'un motif indépendamment de sa position dans l'image."),
        ("• Hiérarchie Spatiale :", "\n  Des contours simples aux textures, puis aux parties d'objets complexes."),
        ("• Efficacité Échantillonnale :", "\n  Apprend efficacement même avec peu d'exemples d'entraînement grâce à ses contraintes structurelles fortes.")
    ]
    for b_t, r_t in cnn_pts:
        p = tf_cb.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.35

    draw_card(s5, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    h_vit = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.55), Inches(5.6), Inches(0.6))
    set_solid(h_vit, C_BLUE_ACCENT)
    tb_vit_h = s5.shapes.add_textbox(Inches(7.0), Inches(1.65), Inches(5.2), Inches(0.4))
    tb_vit_h.text_frame.paragraphs[0].text = "Vision Transformers — ViT (Self-Attention)"
    tb_vit_h.text_frame.paragraphs[0].font.bold = True; tb_vit_h.text_frame.paragraphs[0].font.size = Pt(13); tb_vit_h.text_frame.paragraphs[0].font.color.rgb = C_WHITE

    tb_vit_b = s5.shapes.add_textbox(Inches(7.05), Inches(2.35), Inches(5.1), Inches(4.2))
    tf_vb = tb_vit_b.text_frame; tf_vb.word_wrap = True
    vit_pts = [
        ("• Découpage en Séquence de Patchs :", "\n  L'image 2D est découpée en patchs 16×16 ou 32×32 projetés linéairement (Tokens)."),
        ("• Self-Attention Multi-Têtes Globale :", "\n  Chaque patch interagit directement avec tous les autres dès la première couche."),
        ("• Absence de Biais Inductif Spatial :", "\n  Le modèle doit apprendre lui-même la structure 2D via les positional embeddings."),
        ("• Comportement « Data-Hungry » :", "\n  Exige d'immenses corpus (ImageNet, JFT-300M) pour structurer ses relations d'attention.")
    ]
    for b_t, r_t in vit_pts:
        p = tf_vb.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 6 : ARCHITECTURE GLOBALE DU SYSTÈME & DATA-FLOW
    # =========================================================================
    s6 = prs.slides.add_slide(blank)
    draw_bg(s6)
    add_header(s6, "4. Architecture Globale du Projet & Data-Flow End-to-End")
    add_footer(s6, 6)

    flow_steps = [
        ("1. Pipeline Data", "CUB-200-2011", [
            "• Téléchargement & intégrité",
            "• Splits stratifiés (seed=42)",
            "• Sous-échantillonnage 10%-100%",
            "• Augmentations PyTorch"
        ]),
        ("2. Entraînement & ML", "PyTorch + timm", [
            "• ResNet-50 pré-entraîné",
            "• ViT timm patch 16 / 32",
            "• ViT scratch custom",
            "• AdamW + Cross-Entropy"
        ]),
        ("3. Tracking & Modèles", "MLflow Registry", [
            "• 12 runs loggés (mlflow.db)",
            "• Hyperparamètres & métriques",
            "• Checkpoints (.pth) dans results/",
            "• Modèles complets packagés"
        ]),
        ("4. Déploiement & App", "FastAPI + Next.js", [
            "• API REST (/predict, /attention)",
            "• Extraction attention réelle",
            "• Frontend réactif côte-à-côte",
            "• Conteneurs Docker Compose"
        ]),
    ]

    for i, (title_f, sub_f, pts_f) in enumerate(flow_steps):
        l = Inches(0.9 + i * 2.95)
        draw_card(s6, l, Inches(1.55), Inches(2.75), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
        
        h_f = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.55), Inches(2.75), Inches(0.75))
        set_solid(h_f, C_NAVY_DARK if i % 2 == 0 else C_BLUE_ACCENT)
        
        tb_hf = s6.shapes.add_textbox(l + Inches(0.1), Inches(1.62), Inches(2.55), Inches(0.65))
        tf_hf = tb_hf.text_frame
        p = tf_hf.paragraphs[0]; p.text = title_f; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER
        p2 = tf_hf.add_paragraph(); p2.text = sub_f; p2.font.size = Pt(9); p2.font.color.rgb = C_BORDER_LIGHT; p2.alignment = PP_ALIGN.CENTER

        tb_bf = s6.shapes.add_textbox(l + Inches(0.15), Inches(2.45), Inches(2.45), Inches(4.1))
        tf_bf = tb_bf.text_frame; tf_bf.word_wrap = True
        for pt in pts_f:
            p = tf_bf.add_paragraph()
            p.text = pt; p.font.size = Pt(10.5); p.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(6)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 7 : STRUCTURE MODULAIRE DU CODEBASE
    # =========================================================================
    s7 = prs.slides.add_slide(blank)
    draw_bg(s7)
    add_header(s7, "5. Organisation Modulaire du Codebase")
    add_footer(s7, 7)

    modules = [
        ("/src/data & /data", "Pipeline Données (Rôle A)", [
            "• dataset.py : Classe CUBDataset PyTorch robuste",
            "• transforms.py : Augmentations faible et forte",
            "• splits.py : Génération des sous-échantillons (10% à 100%)",
            "• metadata.csv : 11 788 images étiquetées et vérifiées"
        ]),
        ("/src/models & /src/train", "Modélisation Deep Learning (Rôle B)", [
            "• vit_custom.py : ViT codé intégralement from scratch",
            "• vit_pretrained.py : Wrapper de modèles timm (patch 16 & 32)",
            "• resnet50.py : Adaptation torchvision pour 200 classes",
            "• trainer.py : Boucle de train, AMP fp16, Cosine Scheduler"
        ]),
        ("/app/backend & /app/frontend", "Application Web Démo (Rôle C)", [
            "• main.py & routers/ : Endpoints asynchrones FastAPI",
            "• attention.py : Hook PyTorch d'extraction d'attention",
            "• pages/ & components/ : Interface réactive Next.js 14",
            "• Dockerfile & docker-compose.yml : Conteneurisation"
        ]),
        ("/results & /report", "Expériences & Restitution Scientifique", [
            "• results/runs/*.pth : 12 Checkpoints entraînés et sauvegardés",
            "• evaluate_all_checkpoints.py : Évaluation globale consolidée",
            "• register_all_to_mlflow.py : Exportation vers MLflow Registry",
            "• report/main.tex : Rapport scientifique complet rédigé en LaTeX"
        ]),
    ]

    for idx, (mod_t, mod_sub, mod_pts) in enumerate(modules):
        row = idx // 2
        col = idx % 2
        l = Inches(0.9 + col * 5.9)
        t = Inches(1.55 + row * 2.6)
        
        draw_card(s7, l, t, Inches(5.6), Inches(2.35), C_WHITE, C_BORDER_LIGHT)
        h_m = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(5.6), Inches(0.45))
        set_solid(h_m, C_NAVY_DARK if col == 0 else C_BLUE_ACCENT)
        
        tb_m = s7.shapes.add_textbox(l + Inches(0.15), t + Inches(0.08), Inches(5.3), Inches(0.35))
        tf_m = tb_m.text_frame
        p = tf_m.paragraphs[0]; p.text = f"{mod_t}  —  {mod_sub}"; p.font.bold = True; p.font.size = Pt(11); p.font.color.rgb = C_WHITE

        tb_mb = s7.shapes.add_textbox(l + Inches(0.2), t + Inches(0.55), Inches(5.2), Inches(1.7))
        tf_mb = tb_mb.text_frame; tf_mb.word_wrap = True
        for pt in mod_pts:
            p = tf_mb.add_paragraph()
            p.text = pt; p.font.size = Pt(10.5); p.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(2)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 8 : PIPELINE DE DONNÉES & SPLITS (RÔLE A)
    # =========================================================================
    s8 = prs.slides.add_slide(blank)
    draw_bg(s8)
    add_header(s8, "6. Pipeline de Données & Splits Stratifiés (Rôle A)")
    add_footer(s8, 8)

    draw_card(s8, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s8, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "SPLIT STRATIFIÉ (SEED = 42)", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s8 = s8.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s8 = tb_s8.text_frame; tf_s8.word_wrap = True

    pts_s8 = [
        ("• Répartition Officielle & Validation :", "\n  - Train officiel : 5 994 images\n  - Test officiel : 5 794 images\n  - Ensemble de Validation (15% du train = 900 imgs)\n  - Train final : 5 094 images (100%)"),
        ("• Sous-ensembles pour l'Ablation Data-Hungry :", "\n  - 10%  ->  509 images (~2.5 images / espèce)\n  - 25%  ->  1 273 images (~6 images / espèce)\n  - 50%  ->  2 547 images (~12 images / espèce)\n  - 100% ->  5 094 images (~25 images / espèce)"),
        ("• Stratification Stricte :", "\n  Toutes les 200 classes restent représentées même à 10%.")
    ]
    for b_t, r_t in pts_s8:
        p = tf_s8.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # Right: Table summary
    draw_card(s8, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s8, Inches(7.05), Inches(1.8), Inches(3.2), Inches(0.32), "VOLUMÉTRIE DES SOUS-SPLITS", C_NAVY_DARK, C_WHITE, 9.5)

    rows_s8 = 6
    cols_s8 = 3
    t_shape_s8 = s8.shapes.add_table(rows_s8, cols_s8, Inches(7.05), Inches(2.35), Inches(5.1), Inches(4.1))
    table_s8 = t_shape_s8.table
    table_s8.columns[0].width = Inches(2.1)
    table_s8.columns[1].width = Inches(1.5)
    table_s8.columns[2].width = Inches(1.5)

    headers_s8 = ["Sous-Ensemble", "Nb Images", "Imgs / Classe"]
    for c_idx, h in enumerate(headers_s8):
        cell = table_s8.cell(0, c_idx)
        cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY_DARK
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.size = Pt(10.5); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

    data_s8 = [
        ["Train 10%", "509", "~2.5"],
        ["Train 25%", "1 273", "~6.3"],
        ["Train 50%", "2 547", "~12.7"],
        ["Train 100%", "5 094", "~25.5"],
        ["Validation", "900", "4.5"],
    ]
    for r_idx, row in enumerate(data_s8):
        for c_idx, val in enumerate(row):
            cell = table_s8.cell(r_idx + 1, c_idx)
            cell.fill.solid(); cell.fill.fore_color.rgb = C_WHITE if r_idx % 2 == 0 else C_BG_PAGE
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(10); p.font.color.rgb = C_NAVY_DARK
            p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 9 : ANALYSE EXPLORATOIRE (EDA) & AUGMENTATIONS (RÔLE A)
    # =========================================================================
    s9 = prs.slides.add_slide(blank)
    draw_bg(s9)
    add_header(s9, "7. Analyse Exploratoire & Augmentations (Rôle A)")
    add_footer(s9, 9)

    draw_card(s9, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s9, Inches(1.15), Inches(1.8), Inches(3.0), Inches(0.32), "TRANSFORMATIONS & EDA", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s9 = s9.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s9 = tb_s9.text_frame; tf_s9.word_wrap = True

    pts_s9 = [
        ("• Analyse Exploratoire (EDA) :", "\n  - Résolution native variable (de 140 à 500 px de large).\n  - Distribution équilibrée (41 à 60 images / classe).\n  - Normalisation standardisée ImageNet (mean & std)."),
        ("• Pipeline d'Augmentation Faible (défaut) :", "\n  - Resize(224×224) + RandomHorizontalFlip(p=0.5)\n  - RandomResizedCrop(scale=0.8-1.0)"),
        ("• Pipeline d'Augmentation Forte :", "\n  - ColorJitter (luminosité, contraste, saturation)\n  - RandomRotation(±15°) + RandomErasing(p=0.2)")
    ]
    for b_t, r_t in pts_s9:
        p = tf_s9.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # Right: Image EDA (sans rectangle d'arrière-plan)
    eda_img = os.path.join(REPO, "data", "data_processed", "figures", "eda_similar_classes.png")
    if os.path.exists(eda_img):
        s9.shapes.add_picture(eda_img, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
        lbl = s9.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(5.6), Inches(0.3))
        lbl.text_frame.paragraphs[0].text = "Famille des moineaux : Silhouette identique, variations subtiles de plumage"
        lbl.text_frame.paragraphs[0].font.size = Pt(9.5); lbl.text_frame.paragraphs[0].font.italic = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_TEXT_MUTED; lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 10 : MODÉLISATION DEEP LEARNING (RÔLE B)
    # =========================================================================
    s10 = prs.slides.add_slide(blank)
    draw_bg(s10)
    add_header(s10, "8. Modélisation Deep Learning — Les 3 Architectures (Rôle B)")
    add_footer(s10, 10)

    arch_details = [
        ("ResNet-50 Pré-entraîné", "23.9M Paramètres", C_NAVY_DARK, [
            "• Backbone ResNet-50 torchvision",
            "• Poids ImageNet-1k",
            "• Bloc Bottleneck avec connexions résiduelles",
            "• Tête linéaire 2048 -> 200 classes"
        ]),
        ("ViT Pré-entraîné (timm)", "21.7M - 22.6M Paramètres", C_BLUE_ACCENT, [
            "• vit_small_patch16_224 (196 tokens)",
            "• vit_small_patch32_224 (49 tokens)",
            "• Pré-entraînement ImageNet-1k",
            "• Tête linéaire 384 -> 200 classes"
        ]),
        ("ViT Custom from Scratch", "11.1M Paramètres", C_NAVY_DARK, [
            "• Implémentation custom complète",
            "• Patch Embedding + [CLS] token",
            "• Positional Embedding appris",
            "• 6 Blocs Transformer (384 dim, 6 heads)"
        ]),
    ]

    for i, (t_a, p_a, c_a, pts_a) in enumerate(arch_details):
        l = Inches(0.9 + i * 3.9)
        draw_card(s10, l, Inches(1.55), Inches(3.7), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
        h_a = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.55), Inches(3.7), Inches(0.7))
        set_solid(h_a, c_a)
        
        tb_ha = s10.shapes.add_textbox(l + Inches(0.1), Inches(1.62), Inches(3.5), Inches(0.6))
        tf_ha = tb_ha.text_frame
        p = tf_ha.paragraphs[0]; p.text = t_a; p.font.bold = True; p.font.size = Pt(11.5); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER
        p2 = tf_ha.add_paragraph(); p2.text = p_a; p2.font.size = Pt(9.5); p2.font.color.rgb = C_BORDER_LIGHT; p2.alignment = PP_ALIGN.CENTER

        tb_ba = s10.shapes.add_textbox(l + Inches(0.2), Inches(2.4), Inches(3.3), Inches(4.1))
        tf_ba = tb_ba.text_frame; tf_ba.word_wrap = True
        for pt in pts_a:
            p = tf_ba.add_paragraph()
            p.text = pt; p.font.size = Pt(10.5); p.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(6)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 11 : PROTOCOLE D'ENTRAÎNEMENT & TRACKING MLflow (RÔLE B)
    # =========================================================================
    s11 = prs.slides.add_slide(blank)
    draw_bg(s11)
    add_header(s11, "9. Protocole d'Entraînement & Traçabilité MLflow (Rôle B)")
    add_footer(s11, 11)

    draw_card(s11, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s11, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "HYPERPARAMÈTRES HARMONISÉS", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s11 = s11.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s11 = tb_s11.text_frame; tf_s11.word_wrap = True

    proto_pts = [
        ("• Optimiseur :", " AdamW (lr = 1e-4, weight decay = 0.01)"),
        ("• Fonction de Perte :", " Cross-Entropy Loss"),
        ("• Batch Size & Résolution :", " 32 images / batch, résolution 224×224"),
        ("• Durée :", " 3 époques (fixée pour contrainte de temps de calcul)"),
        ("• Sélection du Modèle :", " Meilleur checkpoint sauvegardé sur Val Accuracy"),
        ("• Seed Fixe :", " 42 sur l'ensemble des 12 expérimentations")
    ]
    for b_t, r_t in proto_pts:
        p = tf_s11.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # Right: MLflow tracking
    draw_card(s11, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s11, Inches(7.05), Inches(1.8), Inches(3.2), Inches(0.32), "TRAÇABILITÉ AVEC MLFLOW", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_mlf = s11.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_mlf = tb_mlf.text_frame; tf_mlf.word_wrap = True

    mlf_pts = [
        ("• Expérience centralisée :", "\n  Nom : CUB-200-2011_ablation dans mlflow.db"),
        ("• 12 Runs enregistrés :", "\n  Toutes les configurations (10%, 25%, 50%, 100%, patch 16/32, scratch/pré-entraîné)."),
        ("• Métriques enregistrées :", "\n  Top-1 test accuracy, Top-5 test accuracy, test loss, temps d'évaluation, nombre de paramètres."),
        ("• Modèles PyTorch packagés :", "\n  Export direct des meilleurs modèles dans MLflow Model Registry.")
    ]
    for b_t, r_t in mlf_pts:
        p = tf_mlf.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 12 : SYNTHÈSE DES RÉSULTATS — FAITS MARQUANTS (AÉRÉ)
    # =========================================================================
    s12 = prs.slides.add_slide(blank)
    draw_bg(s12)
    add_header(s12, "10. Synthèse des Résultats Expérimentaux (5 794 Images Test)")
    add_footer(s12, 12)

    kpis_s12 = [
        ("71.87 %", "ResNet-50 Top-1", "100% données train"),
        ("93.58 %", "ResNet-50 Top-5", "Précision top-5"),
        ("53.85 %", "ViT Patch32 Top-1", "Pré-entraîné ImageNet"),
        ("1.07 % - 2.00 %", "ViT Scratch Top-1", "Niveau du hasard (~0.5%)"),
    ]
    for i, (val, tit, sub) in enumerate(kpis_s12):
        l = Inches(0.9 + i * 2.95)
        draw_card(s12, l, Inches(1.55), Inches(2.75), Inches(1.4), C_WHITE, C_BORDER_LIGHT)
        bar = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.55), Inches(2.75), Inches(0.04))
        set_solid(bar, C_BLUE_ACCENT if i < 2 else C_NAVY_DARK)
        
        tb = s12.shapes.add_textbox(l + Inches(0.1), Inches(1.65), Inches(2.55), Inches(1.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]; p.text = val; p.font.bold = True; p.font.size = Pt(22); p.font.color.rgb = C_BLUE_ACCENT if i == 0 else C_NAVY_DARK; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = tit; p2.font.bold = True; p2.font.size = Pt(11); p2.font.color.rgb = C_NAVY_DARK; p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph(); p3.text = sub; p3.font.size = Pt(9.5); p3.font.color.rgb = C_TEXT_MUTED; p3.alignment = PP_ALIGN.CENTER

    # Enseignements clés
    draw_card(s12, Inches(0.9), Inches(3.2), Inches(11.533), Inches(3.5), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s12, Inches(1.15), Inches(3.45), Inches(3.2), Inches(0.32), "FAITS MARQUANTS POUR LE JURY", C_NAVY_DARK, C_WHITE, 10)

    tb_insights = s12.shapes.add_textbox(Inches(1.15), Inches(3.95), Inches(11.0), Inches(2.5))
    tf_in = tb_insights.text_frame; tf_in.word_wrap = True

    faits_marquants = [
        ("1. Domination de ResNet-50 :", " Avec 71.87% Top-1 et 93.58% Top-5, le CNN surpasse toutes les variantes grâce à son biais inductif de localité."),
        ("2. Rôle vital du Pré-entraînement :", " Le transfert ImageNet permet au ViT de passer de 1.07% à 53.85% (gain de +52.78 points)."),
        ("3. Sensibilité à la taille de patch :", " En 3 époques, le patch 32 (49 tokens) converge nettement plus vite (53.85%) que le patch 16 (9.94% pour 196 tokens)."),
        ("4. Comportement du ViT scratch :", " Sans pré-entraînement, le ViT ne peut pas apprendre les représentations visuelles sur 5 094 images.")
    ]
    for b_t, r_t in faits_marquants:
        p = tf_in.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(11); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 13 : MATRICE EXPÉRIMENTALE DÉTAILLÉE DES 12 CONFIGURATIONS
    # =========================================================================
    s13 = prs.slides.add_slide(blank)
    draw_bg(s13)
    add_header(s13, "11. Matrice Expérimentale Détaillée des 12 Configurations")
    add_footer(s13, 13)

    draw_card(s13, Inches(0.9), Inches(1.55), Inches(11.533), Inches(5.15), C_WHITE, C_BORDER_LIGHT)

    rows_s13 = 7
    cols_s13 = 5
    t_shape_s13 = s13.shapes.add_table(rows_s13, cols_s13, Inches(1.15), Inches(1.8), Inches(11.0), Inches(4.6))
    table_s13 = t_shape_s13.table
    table_s13.columns[0].width = Inches(4.6)
    table_s13.columns[1].width = Inches(1.6)
    table_s13.columns[2].width = Inches(1.6)
    table_s13.columns[3].width = Inches(1.6)
    table_s13.columns[4].width = Inches(1.6)

    headers_s13 = ["Configuration / Modèle", "Top-1 Accuracy", "Top-5 Accuracy", "Test Loss", "Paramètres"]
    for c_idx, h in enumerate(headers_s13):
        cell = table_s13.cell(0, c_idx)
        cell.fill.solid(); cell.fill.fore_color.rgb = C_NAVY_DARK
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.size = Pt(11.5); p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER

    data_s13 = [
        ["ResNet-50 (Pré-entraîné — 100% données)", "71.87 %", "93.58 %", "1.05", "23.9 M"],
        ["ResNet-50 (Pré-entraîné — 50% données)", "48.74 %", "81.12 %", "2.09", "23.9 M"],
        ["ResNet-50 (Pré-entraîné — 25% données)", "23.47 %", "48.88 %", "3.65", "23.9 M"],
        ["ViT Pré-entraîné (Patch 32 — 100% données)", "53.85 %", "83.50 %", "1.81", "22.6 M"],
        ["ViT Pré-entraîné (Patch 16 — 100% données)", "9.94 %", "25.37 %", "4.65", "21.7 M"],
        ["ViT from Scratch (Patch 16/32 — 100% données)", "1.07 % - 2.00 %", "4.75 % - 8.46 %", "~5.15", "11.1 M"],
    ]
    for r_idx, row in enumerate(data_s13):
        for c_idx, val in enumerate(row):
            cell = table_s13.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if r_idx % 2 == 0 else C_BG_PAGE
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(11); p.font.color.rgb = C_NAVY_DARK
            if c_idx == 0:
                p.alignment = PP_ALIGN.LEFT; p.font.bold = (r_idx in [0, 3])
            else:
                p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 14 : ÉTUDE D'ABLATION 1 — COMPORTEMENT DATA-HUNGRY
    # =========================================================================
    s14 = prs.slides.add_slide(blank)
    draw_bg(s14)
    add_header(s14, "12. Étude d'Ablation : Comportement « Data-Hungry »")
    add_footer(s14, 14)

    draw_card(s14, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s14, Inches(1.15), Inches(1.8), Inches(3.0), Inches(0.32), "ANALYSE DE SENSIBILITÉ", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s14 = s14.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s14 = tb_s14.text_frame; tf_s14.word_wrap = True

    pts_s14 = [
        ("Progression Monotone de ResNet-50 :", "\n  • 10% (509 imgs)   : 5.11 % Top-1\n  • 25% (1 273 imgs) : 23.47 % Top-1\n  • 50% (2 547 imgs) : 48.74 % Top-1\n  • 100% (5 094 imgs) : 71.87 % Top-1\n-> Le CNN extrait une structure exploitable dès 500 images grâce à son biais inductif."),
        ("Stagnation du ViT from Scratch :", "\n  • Reste entre 0.97% et 2.16% indépendamment du volume de données.\n-> Confirmation empirique : sans a priori spatial, le ViT ne peut pas structurer son attention sur 5 094 images.")
    ]
    for b_t, r_t in pts_s14:
        p = tf_s14.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    # Right: Curve Image (sans rectangle d'arrière-plan)
    dh_img = os.path.join(REPO, "report", "figures", "data_hungry_curve.png")
    if os.path.exists(dh_img):
        s14.shapes.add_picture(dh_img, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
        lbl = s14.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(5.6), Inches(0.3))
        lbl.text_frame.paragraphs[0].text = "Courbe Data-Hungry : Accuracy Top-1 vs Fraction des données (%)"
        lbl.text_frame.paragraphs[0].font.size = Pt(9.5); lbl.text_frame.paragraphs[0].font.italic = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_TEXT_MUTED; lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 15 : ÉTUDE D'ABLATION 2 — PRÉ-ENTRAÎNEMENT & PATCH SIZE
    # =========================================================================
    s15 = prs.slides.add_slide(blank)
    draw_bg(s15)
    add_header(s15, "13. Étude d'Ablation : Pré-entraînement & Patch Size")
    add_footer(s15, 15)

    draw_card(s15, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s15, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "ENSEIGNEMENTS ARCHITECTURAUX", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s15 = s15.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s15 = tb_s15.text_frame; tf_s15.word_wrap = True

    pts_s15 = [
        ("1. Pré-entraînement indispensable :", "\n• ViT Scratch (Patch 16) : 1.07 % Top-1\n• ViT Pré-entraîné (Patch 32) : 53.85 % Top-1\n-> Gain spectaculaire de +52.78 points apporté par le transfert ImageNet."),
        ("2. Comparaison Patch 16 vs Patch 32 :", "\n• Patch 32 : 53.85 % Top-1 (83.50 % Top-5)\n• Patch 16 : 9.94 % Top-1 (25.37 % Top-5)\n• Explication : Patch 16 engendre une séquence 4x plus longue (196 tokens vs 49 tokens), nécessitant un learning rate adapté et plus d'époques pour converger.")
    ]
    for b_t, r_t in pts_s15:
        p = tf_s15.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    # Right: Final Comparison Image (sans rectangle d'arrière-plan)
    comp_img = os.path.join(REPO, "results", "figures", "final_model_comparison.png")
    if os.path.exists(comp_img):
        s15.shapes.add_picture(comp_img, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
        lbl = s15.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(5.6), Inches(0.3))
        lbl.text_frame.paragraphs[0].text = "Comparatif final : Top-1 Accuracy (%) des différentes variantes"
        lbl.text_frame.paragraphs[0].font.size = Pt(9.5); lbl.text_frame.paragraphs[0].font.italic = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_TEXT_MUTED; lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 16 : ARCHITECTURE SYSTÈME DE L'APPLICATION DÉMO (RÔLE C)
    # =========================================================================
    s16 = prs.slides.add_slide(blank)
    draw_bg(s16)
    add_header(s16, "14. Architecture Système de l'Application Démo (Rôle C)")
    add_footer(s16, 16)

    app_arch_cards = [
        ("Frontend Démo (Next.js 14)", [
            ("• Interface réactive React 18", True),
            ("• Upload d'image par Drag & Drop", False),
            ("• Comparaison côte-à-côte (ViT vs ResNet)", False),
            ("• Barres de confiance Top-3 dynamiques", False),
            ("• Indicateur d'accord / désaccord", False),
            ("• Superposition heatmap interactive", False)
        ]),
        ("Backend API (FastAPI)", [
            ("• Serveur ASGI Uvicorn haute performance", True),
            ("• Endpoint POST /predict : inférence parallèle", False),
            ("• Endpoint POST /attention : extraction heatmap", False),
            ("• Endpoint GET /classes : mapping 200 espèces", False),
            ("• Endpoint GET /health : vérification de santé", False),
            ("• Inférence CPU optimisée en mémoire", False)
        ]),
        ("Moteur d'Inférence PyTorch", [
            ("• Modèles pré-chargés au démarrage (lifespan)", True),
            ("• ResNet-50 & ViT patch32 en mode eval()", False),
            ("• Pipeline de transformation temps réel", False),
            ("• Hook d'attention sur dernier bloc ViT", False),
            ("• Mode commutable via USE_MOCK (dev/prod)", False),
            ("• Sortie normalisée JSON & base64", False)
        ])
    ]

    for i, (titre, points) in enumerate(app_arch_cards):
        l = Inches(0.9 + i * 3.9)
        draw_card(s16, l, Inches(1.55), Inches(3.7), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
        h_s = s16.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(1.55), Inches(3.7), Inches(0.55))
        set_solid(h_s, C_NAVY_DARK if i != 1 else C_BLUE_ACCENT)
        
        tb_h = s16.shapes.add_textbox(l + Inches(0.15), Inches(1.62), Inches(3.4), Inches(0.4))
        tb_h.text_frame.paragraphs[0].text = titre
        tb_h.text_frame.paragraphs[0].font.bold = True
        tb_h.text_frame.paragraphs[0].font.size = Pt(12)
        tb_h.text_frame.paragraphs[0].font.color.rgb = C_WHITE

        tb_b = s16.shapes.add_textbox(l + Inches(0.2), Inches(2.25), Inches(3.3), Inches(4.3))
        tf_b = tb_b.text_frame; tf_b.word_wrap = True
        for j, (txt_pt, is_bold) in enumerate(points):
            p = tf_b.paragraphs[0] if j == 0 else tf_b.add_paragraph()
            p.text = txt_pt; p.font.size = Pt(10.5)
            p.font.bold = is_bold; p.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(4)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 17 : INTERPRÉTABILITÉ — EXTRACTION DE L'ATTENTION ViT (RÔLE C)
    # =========================================================================
    s17 = prs.slides.add_slide(blank)
    draw_bg(s17)
    add_header(s17, "15. Interprétabilité : Carte d'Attention du ViT (Rôle C)")
    add_footer(s17, 17)

    draw_card(s17, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s17, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "MÉCANISME D'EXTRACTION", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s17 = s17.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s17 = tb_s17.text_frame; tf_s17.word_wrap = True

    pts_s17 = [
        ("• Problématique timm :", "\n  timm active par défaut fused_attn=True (scaled_dot_product_attention), ce qui désactive le calcul explicite des poids d'attention."),
        ("• Solution technique (Hook PyTorch) :", "\n  1. Forcer fused_attn=False sur le dernier bloc.\n  2. Intercepter l'entrée de attn_drop via register_forward_hook().\n  3. Moyennage sur les têtes d'attention (forme [B, N, N]).\n  4. Extraction de la ligne du token [CLS] vers les patchs.\n  5. Reshape en grille 2D carrée (ex. 7×7 pour patch 32)."),
        ("• Rendu Heatmap :", "\n  Normalisation min-max, interpolation bilinéaire et superposition RGBA à 50% d'opacité.")
    ]
    for b_t, r_t in pts_s17:
        p = tf_s17.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # Right: Nominal App Image (sans rectangle d'arrière-plan)
    nom_img = os.path.join(REPO, "report", "figures", "app_nominal.png")
    if os.path.exists(nom_img):
        s17.shapes.add_picture(nom_img, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
        lbl = s17.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(5.6), Inches(0.3))
        lbl.text_frame.paragraphs[0].text = "Cas Nominal : Attention focalisée sur les zones anatomiques discriminantes"
        lbl.text_frame.paragraphs[0].font.size = Pt(9.5); lbl.text_frame.paragraphs[0].font.italic = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_TEXT_MUTED; lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 18 : VALIDATION QUALITATIVE & HORS-DISTRIBUTION (OOD)
    # =========================================================================
    s18 = prs.slides.add_slide(blank)
    draw_bg(s18)
    add_header(s18, "16. Robustesse & Détection Hors-Distribution (OOD)")
    add_footer(s18, 18)

    draw_card(s18, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s18, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "COMPORTEMENT EN PRODUCTION", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s18 = s18.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s18 = tb_s18.text_frame; tf_s18.word_wrap = True

    pts_s18 = [
        ("Cas Nominal (Image de Test CUB) :", "\n• Espèce : Gray-crowned Rosy Finch\n• Résultat : Accord total entre ViT et ResNet-50\n• Confiances : ResNet (83.4%), ViT (26.5%)"),
        ("Cas Hors-Distribution (OOD) :", "\n• Espèce testée : Rouge-gorge européen (absent de CUB-200)\n• Résultat : Baisse drastique de la confiance et désaccord entre les modèles.\n• Bénéfice : Signal d'incertitude explicite évitant une erreur silencieuse à forte confiance.")
    ]
    for b_t, r_t in pts_s18:
        p = tf_s18.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11.5); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(10)
        p.line_spacing = 1.35

    # Right: OOD Image (sans rectangle d'arrière-plan)
    ood_img = os.path.join(REPO, "report", "figures", "app_ood.png")
    if os.path.exists(ood_img):
        s18.shapes.add_picture(ood_img, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.8))
        lbl = s18.shapes.add_textbox(Inches(6.8), Inches(6.45), Inches(5.6), Inches(0.3))
        lbl.text_frame.paragraphs[0].text = "Cas Hors-Distribution : Baisse de confiance et désaccord sain entre modèles"
        lbl.text_frame.paragraphs[0].font.size = Pt(9.5); lbl.text_frame.paragraphs[0].font.italic = True
        lbl.text_frame.paragraphs[0].font.color.rgb = C_TEXT_MUTED; lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 19 : INDUSTRIALISATION, DOCKERISATION & DÉPLOIEMENT
    # =========================================================================
    s19 = prs.slides.add_slide(blank)
    draw_bg(s19)
    add_header(s19, "17. Déploiement & Dockerisation Reproductible")
    add_footer(s19, 19)

    docker_grid = [
        ("Backend FastAPI (PyTorch CPU)", [
            "• Image python:3.11-slim optimisée",
            "• PyTorch CPU via --extra-index-url (léger)",
            "• Utilisateur non-root sécurisé",
            "• Healthcheck HTTP sur /health"
        ]),
        ("Frontend Next.js 14 Standalone", [
            "• Build multi-stage Node 20 Alpine",
            "• Support libc6-compat (SWC)",
            "• Mode output: 'standalone' ultra-rapide",
            "• Port 3000 avec injection d'API URL"
        ]),
        ("Orchestration Docker Compose", [
            "• Réseau bridge dédié (vit-cnn-net)",
            "• Frontend synchronisé avec le backend",
            "• Profil MLflow UI optionnel (port 5000)",
            "• Démarrage en 1 commande : compose up -d"
        ]),
        ("Garantie Multi-OS & Zéro Conflit", [
            "• Fonctionne à l'identique sur Windows, Mac & Linux",
            "• Variables d'environnement via .env.example",
            "• Volumes en lecture seule pour les poids .pth",
            "• Élimination des conflits d'environnement"
        ]),
    ]

    for idx, (titre_b, pts_b) in enumerate(docker_grid):
        row = idx // 2
        col = idx % 2
        l = Inches(0.9 + col * 5.9)
        t = Inches(1.55 + row * 2.6)
        
        draw_card(s19, l, t, Inches(5.6), Inches(2.35), C_WHITE, C_BORDER_LIGHT)
        h_bar = s19.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(5.6), Inches(0.45))
        set_solid(h_bar, C_NAVY_DARK if col == 0 else C_BLUE_ACCENT)
        
        tb_h = s19.shapes.add_textbox(l + Inches(0.15), t + Inches(0.08), Inches(5.3), Inches(0.35))
        tb_h.text_frame.paragraphs[0].text = titre_b
        tb_h.text_frame.paragraphs[0].font.bold = True
        tb_h.text_frame.paragraphs[0].font.size = Pt(11.5)
        tb_h.text_frame.paragraphs[0].font.color.rgb = C_WHITE

        tb_b = s19.shapes.add_textbox(l + Inches(0.2), t + Inches(0.55), Inches(5.2), Inches(1.7))
        tf_b = tb_b.text_frame; tf_b.word_wrap = True
        for pt in pts_b:
            p = tf_b.add_paragraph()
            p.text = pt; p.font.size = Pt(10.5); p.font.color.rgb = C_NAVY_DARK
            p.space_after = Pt(2)
            p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 20 : DISCUSSION CRITIQUE & LIMITES DU PROTOCOLE
    # =========================================================================
    s20 = prs.slides.add_slide(blank)
    draw_bg(s20)
    add_header(s20, "18. Discussion Critique & Limites du Protocole")
    add_footer(s20, 20)

    draw_card(s20, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s20, Inches(1.15), Inches(1.8), Inches(3.0), Inches(0.32), "LIMITES DU PROTOCOLE", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s20_l = s20.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s20_l = tb_s20_l.text_frame; tf_s20_l.word_wrap = True

    limites_s20 = [
        ("• Contrainte de 3 époques :", "\n  Temps de calcul restreint qui pénalise le ViT scratch (sous-entraînement manifeste)."),
        ("• Écart Patch 16 vs Patch 32 :", "\n  La variante patch 16 nécessite un réglage fin du learning rate et un temps de convergence plus long."),
        ("• Échelle réduite du modèle scratch :", "\n  Notre ViT custom (11.1M params) reste plus petit que ViT-Base (86M params)."),
        ("• Absence de variance statistique :", "\n  Un seul seed (seed=42) évalué par configuration.")
    ]
    for b_t, r_t in limites_s20:
        p = tf_s20_l.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # Right: Analyse des enseignements
    draw_card(s20, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s20, Inches(7.05), Inches(1.8), Inches(3.2), Inches(0.32), "ENSEIGNEMENTS SCIENTIFIQUES", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s20_r = s20.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s20_r = tb_s20_r.text_frame; tf_s20_r.word_wrap = True

    enseignements_s20 = [
        ("• Biais inductif = Atout décisif :", "\n  Sur petit dataset, le CNN reste mathématiquement mieux adapté pour extraire des textures locales sans pré-entraînement."),
        ("• ViT = Puissance conditionnée :", "\n  Le ViT surpasse les architectures classiques uniquement lorsqu'il est alimenté par un transfert de connaissances massif."),
        ("• Fine-grained = Dépendance au pré-entraînement :", "\n  La tâche fine-grained exige une séparation de classes si fine que les poids ImageNet deviennent indispensables.")
    ]
    for b_t, r_t in enseignements_s20:
        p = tf_s20_r.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 21 : PERSPECTIVES D'AVENIR & ÉVOLUTIONS TECHNIQUES
    # =========================================================================
    s21 = prs.slides.add_slide(blank)
    draw_bg(s21)
    add_header(s21, "19. Perspectives d'Avenir & Évolutions Techniques")
    add_footer(s21, 21)

    draw_card(s21, Inches(0.9), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s21, Inches(1.15), Inches(1.8), Inches(3.2), Inches(0.32), "DÉJÀ CODÉ DANS LE PROJET", C_NAVY_DARK, C_WHITE, 9.5)

    tb_s21_l = s21.shapes.add_textbox(Inches(1.15), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s21_l = tb_s21_l.text_frame; tf_s21_l.word_wrap = True

    persp_dev = [
        ("• Scheduler Cosine Annealing :", "\n  Décroissance cosinusoïdale avec phase de warmup linéaire pour stabiliser l'attention du ViT."),
        ("• Précision Mixte Automatique (AMP fp16) :", "\n  Accélération de l'entraînement et réduction de 50% de l'empreinte mémoire VRAM."),
        ("• Mécanisme d'Early Stopping :", "\n  Arrêt automatique sur stagnation de la validation loss pour éviter le surapprentissage."),
        ("• Script consolidé d'évaluation :", "\n  Réévaluation automatique de tous les checkpoints sur le test set.")
    ]
    for b_t, r_t in persp_dev:
        p = tf_s21_l.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    draw_card(s21, Inches(6.8), Inches(1.55), Inches(5.6), Inches(5.15), C_WHITE, C_BORDER_LIGHT)
    add_pill_badge(s21, Inches(7.05), Inches(1.8), Inches(3.2), Inches(0.32), "PERSPECTIVES DE RECHERCHE", C_BLUE_ACCENT, C_WHITE, 9.5)

    tb_s21_r = s21.shapes.add_textbox(Inches(7.05), Inches(2.3), Inches(5.1), Inches(4.2))
    tf_s21_r = tb_s21_r.text_frame; tf_s21_r.word_wrap = True

    persp_fut = [
        ("• Entraînement Long sur GPU (50 à 100 époques) :", "\n  Donner au ViT from scratch le budget itératif nécessaire pour structurer ses têtes d'attention."),
        ("• Distillation de Connaissances (DeiT) :", "\n  Utiliser ResNet-50 comme réseau professeur (Teacher) pour guider le Transformer (Student)."),
        ("• Part-Based Attention Guidance :", "\n  Forcer les têtes d'attention sur les parties anatomiques de l'oiseau (bec, ailes, yeux)."),
        ("• Intégration de Swin Transformer :", "\n  Tester les fenêtres d'attention glissantes pour réintroduire un biais de localité.")
    ]
    for b_t, r_t in persp_fut:
        p = tf_s21_r.add_paragraph()
        run1 = p.add_run(); run1.text = b_t; run1.font.bold = True; run1.font.size = Pt(11); run1.font.color.rgb = C_NAVY_DARK
        run2 = p.add_run(); run2.text = r_t; run2.font.bold = False; run2.font.size = Pt(10.5); run2.font.color.rgb = C_NAVY_DARK
        p.space_after = Pt(6)
        p.line_spacing = 1.35

    # =========================================================================
    # SLIDE 22 : CONCLUSION, DÉMO EN DIRECT & SESSION DE QUESTIONS
    # =========================================================================
    s22 = prs.slides.add_slide(blank)
    draw_bg(s22)

    r_panel_end = s22.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), 0, Inches(5.033), prs.slide_height)
    set_solid(r_panel_end, C_WHITE)
    r_border_end = s22.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), 0, Inches(0.02), prs.slide_height)
    set_solid(r_border_end, C_BORDER_LIGHT)

    add_pill_badge(s22, Inches(0.9), Inches(0.8), Inches(2.8), Inches(0.35), "SYNTHÈSE DU PROJET", C_BLUE_ACCENT, C_WHITE, 10)

    tb_end = s22.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(7.0), Inches(1.8))
    tf_end = tb_end.text_frame
    p = tf_end.paragraphs[0]; p.text = "Merci pour votre attention !"; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = C_NAVY_DARK
    p2 = tf_end.add_paragraph(); p2.text = "Place aux questions & à la démonstration en direct"; p2.font.size = Pt(15); p2.font.bold = True; p2.font.color.rgb = C_BLUE_ACCENT; p2.space_before = Pt(6)

    takeaways = [
        ("01", "ResNet-50 domine sur petit jeu de données", "Le biais inductif convolutif est capital sur 5 094 images (71.87% Top-1)."),
        ("02", "Le ViT exige un pré-entraînement", "Sans ImageNet, le ViT ne converge pas (+52.78 pts grâce au pré-entraînement)."),
        ("03", "Livrables complets & opérationnels", "Code modulaire, MLflow, API FastAPI, Next.js et conteneurs Docker prêts."),
    ]
    for i, (num, titre_t, desc_t) in enumerate(takeaways):
        top = Inches(3.3 + i * 1.25)
        draw_card(s22, Inches(0.9), top, Inches(7.0), Inches(1.1), C_WHITE, C_BORDER_LIGHT)
        add_diamond_badge(s22, Inches(1.1), top + Inches(0.18), Inches(0.72), num, C_NAVY_DARK if i % 2 == 0 else C_BLUE_ACCENT, C_WHITE)
        
        tb = s22.shapes.add_textbox(Inches(2.05), top + Inches(0.15), Inches(5.6), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = titre_t; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = C_NAVY_DARK
        p2 = tf.add_paragraph(); p2.text = desc_t; p2.font.size = Pt(10); p2.font.color.rgb = C_TEXT_MUTED; p2.line_spacing = 1.35

    # Right: Live Demo Links
    tb_demo_t = s22.shapes.add_textbox(Inches(8.6), Inches(0.8), Inches(4.4), Inches(0.4))
    tb_demo_t.text_frame.paragraphs[0].text = "DÉMONSTRATION EN DIRECT"
    tb_demo_t.text_frame.paragraphs[0].font.bold = True; tb_demo_t.text_frame.paragraphs[0].font.size = Pt(13); tb_demo_t.text_frame.paragraphs[0].font.color.rgb = C_NAVY_DARK

    demo_cards = [
        ("Frontend Web (Next.js 14)", "http://localhost:3000", "Interface réactive avec drag & drop d'oiseaux, prédictions Top-3 et Heatmap."),
        ("Backend API (FastAPI Docs)", "http://localhost:8000/docs", "Documentation Swagger interactive des endpoints /predict et /attention."),
        ("Suivi MLflow UI", "http://localhost:5000", "Visualisation des métriques des 12 runs et des modèles packagés.")
    ]
    for i, (title_d, url_d, info_d) in enumerate(demo_cards):
        top = Inches(1.4 + i * 1.8)
        draw_card(s22, Inches(8.6), top, Inches(4.4), Inches(1.6), C_WHITE, C_BORDER_LIGHT)
        
        tb = s22.shapes.add_textbox(Inches(8.85), top + Inches(0.12), Inches(4.0), Inches(1.35))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title_d; p.font.bold = True; p.font.size = Pt(12); p.font.color.rgb = C_NAVY_DARK
        p2 = tf.add_paragraph(); p2.text = url_d; p2.font.bold = True; p2.font.size = Pt(10.5); p2.font.color.rgb = C_BLUE_ACCENT; p2.space_after = Pt(2)
        p3 = tf.add_paragraph(); p3.text = info_d; p3.font.size = Pt(9.5); p3.font.color.rgb = C_TEXT_MUTED; p3.line_spacing = 1.35

    # ── Notes d'orateur complètes intégrées par slide ───
    s1.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "Bonjour Monsieur le Professeur, bonjour aux membres du jury. "
        "Nous vous présentons notre projet annuel de Master 1 IA : 'Vision Transformer vs CNN : une étude comparative sur la classification fine-grained d'oiseaux'. "
        "Notre équipe est composée d'Abdoul-raouf (Rôle A, Data), Billy (Rôle B, Modèles & MLflow) et Koffi (Rôle C, Backend, Frontend & Docker). "
        "Nous avons entraîné et évalué 12 configurations expérimentales complètes sur CUB-200-2011 et livré une application web dockerisée."
    )
    s2.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "Pour vous présenter notre travail, nous avons structuré notre soutenance en 6 grandes étapes : "
        "le contexte et la théorie, l'architecture globale du système, le pipeline data, la modélisation et les résultats expérimentaux, "
        "les études d'ablation approfondies, et enfin l'application web, la dockerisation et nos perspectives critiques."
    )
    s3.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "La classification fine-grained consiste à distinguer 200 espèces d'oiseaux très proches. "
        "La variabilité inter-classe est minime (nuances subtiles de plumage ou de bec) et la variabilité intra-classe est forte (poses, éclairages). "
        "Le benchmark CUB-200-2011 compte 11 788 images, ce qui représente un volume modeste pour le Deep Learning."
    )
    s4.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "Le benchmark CUB-200 est notre terrain d'expérimentation standardisé avec 200 classes et 5 794 images de test. "
        "Nos objectifs de recherche sont triples : comparer objectivement CNN et ViT, valider empiriquement la sensibilité au volume de données (Data-Hungry), "
        "et fournir une application web explicable en temps réel."
    )
    s5.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "Deux paradigmes s'opposent : le CNN (ResNet-50) possède un biais inductif de localité spatiale qui le rend très efficace sur petit dataset. "
        "Le Vision Transformer (ViT) découpe l'image en patchs et utilise la Self-Attention globale sans a priori spatial. "
        "L'hypothèse théorique est que le ViT est data-hungry. Je passe la parole à Billy pour l'architecture et les modèles."
    )
    s6.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Merci Abdoul-raouf. Voici la vue d'ensemble de notre chaîne end-to-end : "
        "du pipeline de données stratifié à l'entraînement PyTorch, au tracking centralisé dans MLflow, "
        "jusqu'au déploiement dans l'API FastAPI et l'interface Next.js conteneurisées sous Docker."
    )
    s7.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Le codebase est strictement modulaire : /src/data pour le dataset et les transforms d'Abdoul-raouf, "
        "/src/models et /src/train pour les réseaux et la boucle d'entraînement, /app pour l'API et le frontend de Koffi, "
        "et /results pour les 12 checkpoints .pth et métriques évaluées."
    )
    s8.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "Pour le rôle A, j'ai extrait un split de validation stratifié de 15% (900 images) avec un seed fixe (42). "
        "J'ai ensuite créé quatre sous-échantillons d'entraînement : 10% (509 imgs), 25% (1273 imgs), 50% (2547 imgs) et 100% (5094 imgs). "
        "La stratification garantit que les 200 espèces restent toutes représentées même à 10%."
    )
    s9.notes_slide.notes_text_frame.text = (
        "[ABDOUL-RAOUF — RÔLE A]\n"
        "L'analyse exploratoire a montré des résolutions très variables, d'où un resize à 224x224 normalisé ImageNet. "
        "La figure montre la ressemblance extrême entre moineaux. "
        "J'ai mis en place deux pipelines d'augmentation : un faible (défaut) et un fort avec ColorJitter, rotations et RandomErasing. Billy va vous présenter les modèles."
    )
    s10.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Nous avons conçu 3 architectures adaptées aux 200 classes : ResNet-50 pré-entraîné (23.9M params), "
        "ViT pré-entraîné timm small patch 16 et 32 (21.7M - 22.6M params), et un ViT custom codé from scratch (11.1M params) "
        "avec 6 blocs Transformer, dimension 384 et 6 têtes d'attention, sans aucun pré-entraînement."
    )
    s11.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Le protocole est harmonisé : optimiseur AdamW (lr 1e-4, weight decay 0.01), Cross-Entropy, batch de 32 sur 3 époques. "
        "Le meilleur modèle est sauvé sur la validation accuracy. "
        "Toutes les 12 configurations ont été loggées dans MLflow pour tracer les courbes, métriques et artefacts."
    )
    s12.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Sur les 5 794 images de test : ResNet-50 pré-entraîné 100% réalise 71.87% en Top-1 et 93.58% en Top-5. "
        "Le ViT pré-entraîné patch 32 atteint 53.85% (83.50% top-5). "
        "Le ViT from scratch reste bloqué entre 1.07% et 2.00% (proche du hasard à 0.5%)."
    )
    s13.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Voici la matrice complète de nos 12 expérimentations. Vous pouvez constater la cohérence mathématique : "
        "la précision de ResNet-50 s'élève de manière strictement monotone avec la volumétrie de données."
    )
    s14.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "C'est le résultat central : ResNet-50 progresse de façon strictement monotone (5.11% à 10%, 23.47% à 25%, 48.74% à 50% et 71.87% à 100%). "
        "Le ViT scratch reste plat à ~1%. Sans biais inductif, le ViT ne peut pas apprendre de représentations sur 5 000 images."
    )
    s15.notes_slide.notes_text_frame.text = (
        "[BILLY — RÔLE B]\n"
        "Deux conclusions : le pré-entraînement ImageNet fait bondir le ViT de +52.78 points. "
        "Sur 3 époques, le patch 32 (53.85%) surpasse le patch 16 (9.94%) car sa séquence de 49 tokens s'adapte beaucoup plus vite que les 196 tokens du patch 16. "
        "Je passe la parole à Koffi pour l'application et le déploiement."
    )
    s16.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "Merci Billy. Pour le rôle C, j'ai développé l'application fullstack : "
        "un frontend réactif Next.js 14 en React 18 avec upload drag & drop et comparaison côte-à-côte, "
        "un backend FastAPI asynchrone (/predict, /attention, /classes), et un moteur PyTorch avec modèles pré-chargés en mémoire vive au démarrage."
    )
    s17.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "Pour l'interprétabilité, nous avons développé un hook PyTorch sur attn_drop du dernier bloc Transformer. "
        "On extrait les poids d'attention du token CLS vers les patchs, moyennés sur les têtes, reshapés en grille 2D et superposés en heatmap. "
        "L'attention se focalise précisément sur la tête, le bec et les ailes de l'oiseau."
    )
    s18.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "En test hors-distribution avec un Rouge-gorge européen (absent de CUB), "
        "les confiances s'effondrent sous les 10% et les deux modèles sont en désaccord. "
        "C'est un comportement très sain : le système émet un signal explicite d'incertitude plutôt qu'une fausse certitude."
    )
    s19.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "Le projet est 100% dockerisé : backend python:3.11-slim optimisé PyTorch CPU sans surcoût CUDA, "
        "frontend multi-stage standalone Node 20 Alpine, et Docker Compose avec healthchecks synchronisés. "
        "Tout démarre en une seule commande : docker compose up -d."
    )
    s20.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "En recul critique : le budget de 3 époques pénalise le ViT scratch (sous-entraînement), "
        "le patch 16 a besoin d'un learning rate adapté et d'un temps plus long, et notre ViT custom (11.1M) est plus petit que ViT-Base. "
        "Néanmoins, les tendances corroborent parfaitement la littérature scientifique."
    )
    s21.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C]\n"
        "Nous avons déjà codé dans le repo : le Scheduler Cosine avec Warmup, la Précision Mixte AMP fp16, et l'Early Stopping. "
        "Les perspectives futures sont un entraînement long sur GPU (50-100 époques), la distillation de connaissances DeiT, et l'attention guidée par parties anatomiques."
    )
    s22.notes_slide.notes_text_frame.text = (
        "[KOFFI — RÔLE C & TOUTE L'ÉQUIPE]\n"
        "En conclusion : ResNet-50 domine sur petit dataset (71.87% Top-1), le ViT exige un pré-entraînement (+52.78 pts), "
        "et nous livrons un projet complet, reproductible et dockerisé. "
        "Monsieur le Professeur, nous vous remercions pour votre écoute et nous sommes prêts pour vos questions et la démonstration en direct !"
    )

    # ── Sauvegarde dans le fichier unique officiel ───
    target_path = os.path.join(REPO, "presentation_vit_vs_cnn.pptx")
    prs.save(target_path)
    print(f"[SUCCÈS] Présentation officielle 22 slides générée avec succès : {target_path}")

if __name__ == "__main__":
    create_presentation()
